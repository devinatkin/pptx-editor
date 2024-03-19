import pptx
import argparse

def replace_text(pptx, old_text, new_text):
    """
    Replace all occurrences of old_text with new_text in the given pptx object.
    """
    for slide in pptx.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.text = shape.text_frame.text.replace(old_text, new_text)

def replace_text_w_image(pptx, old_text, img_path):
    """
    Replace all occurrences of old_text with new_text in the given pptx object.
    """
    for slide in pptx.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text_frame.text == old_text:
                    shape.text_frame.clear()  # Clear the existing text
                    shape.text_frame.add_paragraph()  # Add a new paragraph
                    slide.shapes.add_picture(img_path, shape.left, shape.top, shape.width, shape.height)

def process_replacement_dict(pptx, replacement_text_dict, replacement_image_dict):
    """
    Replace images and text according to their respective dictionaries.
    """

    for old_text, new_text in replacement_text_dict.items():
        replace_text(pptx, old_text, new_text)

    for old_text, img_path in replacement_image_dict.items():
        replace_text_w_image(pptx, old_text, img_path)

import ast

def str_to_dict(arg):
    """
    Convert a string representation of a dictionary into an actual dictionary.
    """
    return ast.literal_eval(arg)

if __name__ == '__main__':
    # Parse the command line arguments
    # Example Call
    # python editpptx.py TemplateFile.pptx output.pptx "{'<TEXT1>': 'new_text1', '<TEXT2>': 'new_text2'}" "{'<IMAGE1>': 'leaf_image.jpeg','<IMAGE2>': 'leaf_image.jpeg'}"

    parser = argparse.ArgumentParser(description='PPTX Editor')
    parser.add_argument('input_file', type=str, help='Path to the input PPTX file')
    parser.add_argument('output_file', type=str, help='Path to the output PPTX file')
    parser.add_argument('replacements_text', type=str_to_dict, help='Dictionary of text replacements')
    parser.add_argument('replacements_image', type=str_to_dict, help='Dictionary of image replacements')
    args = parser.parse_args()

    pptx = pptx.Presentation(args.input_file)

    # Perform the replacements
    process_replacement_dict(pptx, args.replacements_text, args.replacements_image)

    pptx.save(args.output_file)